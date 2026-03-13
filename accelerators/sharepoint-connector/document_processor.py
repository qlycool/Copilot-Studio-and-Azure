"""
Document text extraction.
Supports the same formats as the Azure AI Search SharePoint connector:
  PDF, DOCX/DOCM, XLSX/XLSM, PPTX/PPTM, CSV, EML, EPUB, HTML, JSON,
  KML/XML, Markdown, MSG, ODT/ODS/ODP, RTF, TXT, ZIP/GZ.

Old binary formats (DOC, XLS, PPT) are not supported — they require
server-side rendering (LibreOffice, etc.).
"""

import io
import logging
from pathlib import PurePosixPath

logger = logging.getLogger(__name__)

# Safety limits for archive extraction
_MAX_DECOMPRESSED_SIZE = 200 * 1024 * 1024  # 200 MB total extracted content
_MAX_ARCHIVE_DEPTH = 3  # max nesting of archives inside archives


def extract_text(filename: str, content: bytes) -> str:
    """
    Extract plain text from file content based on the file extension.
    Returns empty string if extraction fails or format is unsupported.
    """
    ext = PurePosixPath(filename).suffix.lower()

    extractors = {
        # PDF
        ".pdf": _extract_pdf,
        # Microsoft Office (XML-based)
        ".docx": _extract_docx,
        ".docm": _extract_docx,
        ".xlsx": _extract_xlsx,
        ".xlsm": _extract_xlsx,
        ".pptx": _extract_pptx,
        ".pptm": _extract_pptx,
        ".msg": _extract_msg,
        # Open Document
        ".odt": _extract_odt,
        ".ods": _extract_ods,
        ".odp": _extract_odp,
        # Plain text / data
        ".txt": _extract_text,
        ".md": _extract_text,
        ".csv": _extract_csv,
        ".json": _extract_json,
        ".xml": _extract_xml,
        ".kml": _extract_xml,
        # Markup
        ".html": _extract_html,
        ".htm": _extract_html,
        # Rich text / email / ebook
        ".rtf": _extract_rtf,
        ".eml": _extract_eml,
        ".epub": _extract_epub,
        # Archives (extract and process contents)
        ".zip": _extract_zip,
        ".gz": _extract_gz,
    }

    extractor = extractors.get(ext)
    if not extractor:
        logger.warning(f"Unsupported file format: {ext} ({filename})")
        return ""

    try:
        text = extractor(content, filename)
        logger.debug(f"Extracted {len(text)} chars from {filename}")
        return text
    except Exception as e:
        logger.error(f"Failed to extract text from {filename}: {e}")
        return ""


# ------------------------------------------------------------------
# PDF
# ------------------------------------------------------------------

def _extract_pdf(content: bytes, filename: str = "") -> str:
    """Extract text from PDF using PyMuPDF."""
    import fitz  # PyMuPDF

    text_parts: list[str] = []
    with fitz.open(stream=content, filetype="pdf") as doc:
        for page in doc:
            text_parts.append(page.get_text())
    return "\n".join(text_parts)


# ------------------------------------------------------------------
# Microsoft Office (XML-based)
# ------------------------------------------------------------------

def _extract_docx(content: bytes, filename: str = "") -> str:
    """Extract text from DOCX/DOCM using python-docx."""
    from docx import Document

    doc = Document(io.BytesIO(content))
    paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]
    return "\n".join(paragraphs)


def _extract_xlsx(content: bytes, filename: str = "") -> str:
    """Extract text from XLSX/XLSM using openpyxl."""
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(content), read_only=True, data_only=True)
    text_parts: list[str] = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        text_parts.append(f"--- Sheet: {sheet_name} ---")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            line = "\t".join(cells).strip()
            if line:
                text_parts.append(line)
    wb.close()
    return "\n".join(text_parts)


def _extract_pptx(content: bytes, filename: str = "") -> str:
    """Extract text from PPTX/PPTM using python-pptx."""
    from pptx import Presentation

    prs = Presentation(io.BytesIO(content))
    text_parts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        text_parts.append(text)
    return "\n".join(text_parts)


def _extract_msg(content: bytes, filename: str = "") -> str:
    """Extract text from Outlook MSG files using extract-msg."""
    import extract_msg

    msg = extract_msg.Message(io.BytesIO(content))
    parts = []
    if msg.subject:
        parts.append(f"Subject: {msg.subject}")
    if msg.sender:
        parts.append(f"From: {msg.sender}")
    if msg.date:
        parts.append(f"Date: {msg.date}")
    if msg.body:
        parts.append(msg.body)
    msg.close()
    return "\n".join(parts)


# ------------------------------------------------------------------
# Open Document formats
# ------------------------------------------------------------------

def _extract_odt(content: bytes, filename: str = "") -> str:
    """Extract text from ODT (Open Document Text)."""
    from odf.opendocument import load as odf_load
    from odf import text as odf_text

    doc = odf_load(io.BytesIO(content))
    paragraphs: list[str] = []
    for p in doc.getElementsByType(odf_text.P):
        text = _odf_element_text(p)
        if text.strip():
            paragraphs.append(text.strip())
    return "\n".join(paragraphs)


def _extract_ods(content: bytes, filename: str = "") -> str:
    """Extract text from ODS (Open Document Spreadsheet)."""
    from odf.opendocument import load as odf_load
    from odf import table as odf_table, text as odf_text

    doc = odf_load(io.BytesIO(content))
    text_parts: list[str] = []
    for sheet in doc.getElementsByType(odf_table.Table):
        sheet_name = sheet.getAttribute("name") or "Sheet"
        text_parts.append(f"--- Sheet: {sheet_name} ---")
        for row in sheet.getElementsByType(odf_table.TableRow):
            cells = []
            for cell in row.getElementsByType(odf_table.TableCell):
                cell_text = ""
                for p in cell.getElementsByType(odf_text.P):
                    cell_text += _odf_element_text(p)
                cells.append(cell_text)
            line = "\t".join(cells).strip()
            if line:
                text_parts.append(line)
    return "\n".join(text_parts)


def _extract_odp(content: bytes, filename: str = "") -> str:
    """Extract text from ODP (Open Document Presentation)."""
    from odf.opendocument import load as odf_load
    from odf import text as odf_text

    doc = odf_load(io.BytesIO(content))
    paragraphs: list[str] = []
    for p in doc.getElementsByType(odf_text.P):
        text = _odf_element_text(p)
        if text.strip():
            paragraphs.append(text.strip())
    return "\n".join(paragraphs)


def _odf_element_text(element) -> str:
    """Recursively extract text from an ODF element."""
    text = ""
    if hasattr(element, "data"):
        text += element.data
    for child in element.childNodes:
        text += _odf_element_text(child)
    return text


# ------------------------------------------------------------------
# Plain text / data formats
# ------------------------------------------------------------------

def _extract_text(content: bytes, filename: str = "") -> str:
    """Extract text from plain text files."""
    for encoding in ("utf-8", "utf-8-sig", "cp1252", "latin-1"):
        try:
            return content.decode(encoding)
        except (UnicodeDecodeError, ValueError):
            continue
    return content.decode("utf-8", errors="replace")


def _extract_csv(content: bytes, filename: str = "") -> str:
    """Extract text from CSV, preserving rows as tab-separated lines."""
    import csv

    text = _extract_text(content)
    reader = csv.reader(io.StringIO(text))
    lines = ["\t".join(row) for row in reader if any(cell.strip() for cell in row)]
    return "\n".join(lines)


def _extract_json(content: bytes, filename: str = "") -> str:
    """Extract text from JSON by pretty-printing the structure."""
    import json

    text = _extract_text(content)
    try:
        data = json.loads(text)
        return json.dumps(data, indent=2, ensure_ascii=False)
    except json.JSONDecodeError:
        return text  # Return raw text if not valid JSON


def _extract_xml(content: bytes, filename: str = "") -> str:
    """Extract text content from XML/KML files."""
    import xml.etree.ElementTree as ET

    text = _extract_text(content)
    try:
        root = ET.fromstring(text)
        parts: list[str] = []
        for elem in root.iter():
            if elem.text and elem.text.strip():
                parts.append(elem.text.strip())
            if elem.tail and elem.tail.strip():
                parts.append(elem.tail.strip())
        return "\n".join(parts)
    except ET.ParseError:
        return text  # Fallback to raw text


# ------------------------------------------------------------------
# Markup / rich text
# ------------------------------------------------------------------

def _extract_html(content: bytes, filename: str = "") -> str:
    """Extract text from HTML using BeautifulSoup."""
    from bs4 import BeautifulSoup

    text = _extract_text(content)
    soup = BeautifulSoup(text, "html.parser")

    # Remove script and style elements
    for element in soup(["script", "style", "nav", "footer", "header"]):
        element.decompose()

    return soup.get_text(separator="\n", strip=True)


def _extract_rtf(content: bytes, filename: str = "") -> str:
    """Extract text from RTF using striprtf."""
    from striprtf.striprtf import rtf_to_text

    text = _extract_text(content)
    return rtf_to_text(text)


# ------------------------------------------------------------------
# Email / ebook
# ------------------------------------------------------------------

def _extract_eml(content: bytes, filename: str = "") -> str:
    """Extract text from EML (email) files using stdlib email module."""
    import email
    from email import policy

    msg = email.message_from_bytes(content, policy=policy.default)
    parts = []
    if msg["subject"]:
        parts.append(f"Subject: {msg['subject']}")
    if msg["from"]:
        parts.append(f"From: {msg['from']}")
    if msg["date"]:
        parts.append(f"Date: {msg['date']}")

    body = msg.get_body(preferencelist=("plain", "html"))
    if body:
        body_content = body.get_content()
        if body.get_content_type() == "text/html":
            body_content = _extract_html(body_content.encode("utf-8"))
        parts.append(body_content)

    return "\n".join(parts)


def _extract_epub(content: bytes, filename: str = "") -> str:
    """Extract text from EPUB using ebooklib + BeautifulSoup."""
    import ebooklib
    from ebooklib import epub
    from bs4 import BeautifulSoup

    book = epub.read_epub(io.BytesIO(content))
    text_parts: list[str] = []
    for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT):
        soup = BeautifulSoup(item.get_content(), "html.parser")
        text = soup.get_text(separator="\n", strip=True)
        if text:
            text_parts.append(text)
    return "\n".join(text_parts)


# ------------------------------------------------------------------
# Archives
# ------------------------------------------------------------------

def _extract_zip(content: bytes, filename: str = "", _depth: int = 0) -> str:
    """Extract and process files inside a ZIP archive with safety limits."""
    import zipfile

    if _depth >= _MAX_ARCHIVE_DEPTH:
        logger.warning(f"Archive nesting too deep ({_depth}), skipping: {filename}")
        return ""

    text_parts: list[str] = []
    total_extracted = 0
    with zipfile.ZipFile(io.BytesIO(content)) as zf:
        for info in zf.infolist():
            if info.is_dir():
                continue
            # Check decompressed size before extracting
            if info.file_size > _MAX_DECOMPRESSED_SIZE:
                logger.warning(f"File too large inside ZIP: {info.filename} ({info.file_size} bytes), skipping")
                continue
            total_extracted += info.file_size
            if total_extracted > _MAX_DECOMPRESSED_SIZE:
                logger.warning(f"Total extracted size exceeds limit in {filename}, stopping")
                break
            inner_name = info.filename
            try:
                inner_content = zf.read(inner_name)
                # Pass depth for nested archive detection
                inner_ext = PurePosixPath(inner_name).suffix.lower()
                if inner_ext in (".zip", ".gz"):
                    if inner_ext == ".zip":
                        inner_text = _extract_zip(inner_content, inner_name, _depth + 1)
                    else:
                        inner_text = _extract_gz(inner_content, inner_name, _depth + 1)
                else:
                    inner_text = extract_text(inner_name, inner_content)
                if inner_text:
                    text_parts.append(f"--- {inner_name} ---")
                    text_parts.append(inner_text)
            except Exception as e:
                logger.warning(f"Failed to extract {inner_name} from {filename}: {e}")
    return "\n".join(text_parts)


def _extract_gz(content: bytes, filename: str = "", _depth: int = 0) -> str:
    """Extract and process the file inside a GZ archive."""
    import gzip

    if _depth >= _MAX_ARCHIVE_DEPTH:
        logger.warning(f"Archive nesting too deep ({_depth}), skipping: {filename}")
        return ""

    try:
        decompressed = gzip.decompress(content)
    except Exception as e:
        logger.error(f"Failed to decompress {filename}: {e}")
        return ""

    if len(decompressed) > _MAX_DECOMPRESSED_SIZE:
        logger.warning(f"Decompressed size exceeds limit for {filename} ({len(decompressed)} bytes), skipping")
        return ""

    # Strip .gz to get the inner filename
    inner_name = filename.rsplit(".gz", 1)[0] if filename.endswith(".gz") else filename
    if not inner_name or inner_name == filename:
        inner_name = "file.txt"

    inner_ext = PurePosixPath(inner_name).suffix.lower()
    if inner_ext in (".zip", ".gz"):
        if inner_ext == ".zip":
            return _extract_zip(decompressed, inner_name, _depth + 1)
        else:
            return _extract_gz(decompressed, inner_name, _depth + 1)

    return extract_text(inner_name, decompressed)
